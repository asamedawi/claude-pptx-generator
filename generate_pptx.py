import argparse
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).parent
TEMPLATE_PPTX = ROOT / "templates" / "template.pptx"
OUTPUT_PPTX = ROOT / "output" / "output.pptx"
MARKDOWN_LINK_PATTERN = re.compile(r"\[([^\]]+)\]\((https?://[^\s)]+)\)")
RAW_URL_PATTERN = re.compile(r"https?://[^\s]+")


def extract_content_section(text: str) -> list[str]:
    content = text.strip()
    if not content:
        return []
    return [line.rstrip() for line in content.splitlines()]


def finalize_slide(slide_data: dict[str, list[str] | str]) -> dict[str, list[str] | str]:
    bullets = slide_data["bullets"]
    bullet_items = bullets if isinstance(bullets, list) else [str(bullets)]
    slide_data["kind"] = "section" if not bullet_items else "content"
    return slide_data


def parse_slides(lines: list[str]) -> tuple[str, list[dict[str, list[str] | str]]]:
    title = "資料"
    slides: list[dict[str, list[str] | str]] = []
    current_slide: dict[str, list[str] | str] | None = None

    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue
        if line == "ここに資料化したい内容を書く。":
            continue
        if line.startswith("# "):
            title = line[2:].strip() or title
            continue
        if line.startswith("## "):
            if current_slide:
                slides.append(finalize_slide(current_slide))
            current_slide = {"title": line[3:].strip() or "無題", "bullets": []}
            continue
        if line.startswith("- "):
            if current_slide is None:
                current_slide = {"title": "内容", "bullets": []}
            bullets = current_slide["bullets"]
            if isinstance(bullets, list):
                bullets.append(line[2:].strip())
            continue
        if current_slide is None:
            current_slide = {"title": "内容", "bullets": []}
        bullets = current_slide["bullets"]
        if isinstance(bullets, list):
            bullets.append(line)

    if current_slide:
        slides.append(finalize_slide(current_slide))

    if not slides:
        slides = [{"title": "内容", "bullets": ["Markdown を指定するとスライドを生成できます。"], "kind": "content"}]

    return title, slides


def set_shape_text(shape, text: str) -> None:
    if hasattr(shape, "text"):
        shape.text = text


def add_textbox(slide, text: str, left: float, top: float, width: float, height: float, font_size: int) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    if paragraph.runs:
        paragraph.runs[0].font.size = Inches(font_size / 72)


def split_trailing_punctuation(url: str) -> tuple[str, str]:
    trailing_chars = ".,);:!?]}>\"'」』】"
    trimmed_url = url.rstrip(trailing_chars)
    trailing = url[len(trimmed_url) :]
    return trimmed_url, trailing


def add_run(paragraph, text: str, hyperlink: str | None = None) -> None:
    if not text:
        return
    run = paragraph.add_run()
    run.text = text
    if hyperlink:
        run.hyperlink.address = hyperlink


def write_text_with_links(paragraph, text: str) -> None:
    position = 0

    while position < len(text):
        markdown_match = MARKDOWN_LINK_PATTERN.search(text, position)
        raw_url_match = RAW_URL_PATTERN.search(text, position)
        matches = [match for match in [markdown_match, raw_url_match] if match is not None]

        if not matches:
            add_run(paragraph, text[position:])
            break

        match = min(matches, key=lambda item: item.start())

        if match.start() > position:
            add_run(paragraph, text[position:match.start()])

        if match.re == MARKDOWN_LINK_PATTERN:
            add_run(paragraph, match.group(1), match.group(2))
            position = match.end()
            continue

        raw_url = match.group(0)
        clean_url, trailing = split_trailing_punctuation(raw_url)
        if clean_url:
            add_run(paragraph, clean_url, clean_url)
        if trailing:
            add_run(paragraph, trailing)
        position = match.end()


def find_layout(prs: Presentation, preferred_names: list[str], fallback_index: int = 0):
    for name in preferred_names:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
    if len(prs.slide_layouts) > fallback_index:
        return prs.slide_layouts[fallback_index]
    return prs.slide_layouts[0]


def find_body_placeholder(slide):
    for placeholder in slide.placeholders:
        if placeholder == slide.shapes.title:
            continue
        if hasattr(placeholder, "text_frame"):
            return placeholder
    return None


def add_title_slide(prs: Presentation, title: str) -> None:
    layout = find_layout(prs, ["表紙"], 0)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        set_shape_text(slide.shapes.title, title)
    else:
        add_textbox(slide, title, 1, 1.2, 8, 1, 28)

    subtitle = find_body_placeholder(slide)
    if subtitle is not None:
        set_shape_text(subtitle, "Generated by claude-pptx-generator")


def add_section_slide(prs: Presentation, slide_data: dict[str, list[str] | str]) -> None:
    layout = find_layout(prs, ["中表紙", "中々表紙"], 1)
    slide = prs.slides.add_slide(layout)
    title = str(slide_data["title"])

    if slide.shapes.title is not None:
        set_shape_text(slide.shapes.title, title)
    else:
        add_textbox(slide, title, 1, 1.2, 8, 1, 24)

    subtitle = find_body_placeholder(slide)
    if subtitle is not None:
        set_shape_text(subtitle, "")


def add_content_slide(prs: Presentation, slide_data: dict[str, list[str] | str]) -> None:
    layout = find_layout(prs, ["大見出しとコンテンツ", "レイアウト"], 3)
    slide = prs.slides.add_slide(layout)

    title = str(slide_data["title"])
    bullets = slide_data["bullets"]

    if slide.shapes.title is not None:
        set_shape_text(slide.shapes.title, title)
    else:
        add_textbox(slide, title, 0.8, 0.5, 8.5, 0.8, 24)

    body_placeholder = find_body_placeholder(slide)
    bullet_items = bullets if isinstance(bullets, list) else [str(bullets)]

    if body_placeholder is not None and hasattr(body_placeholder, "text_frame"):
        text_frame = body_placeholder.text_frame
        text_frame.clear()
        for index, bullet in enumerate(bullet_items):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            write_text_with_links(paragraph, str(bullet))
            paragraph.level = 0
    else:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(4))
        text_frame = textbox.text_frame
        for index, bullet in enumerate(bullet_items):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            write_text_with_links(paragraph, str(bullet))


def build_argument_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-file", type=Path)
    parser.add_argument("--source-text")
    parser.add_argument("--output", type=Path, default=OUTPUT_PPTX)
    parser.add_argument("--no-open", action="store_true")
    return parser


def load_source_text(args: argparse.Namespace) -> str:
    if args.source_text:
        return args.source_text
    if args.source_file:
        return args.source_file.read_text(encoding="utf-8")
    if not sys.stdin.isatty():
        content = sys.stdin.read()
        if content.strip():
            return content
    raise ValueError("Markdown を --source-file、--source-text、または標準入力で指定してください。")


def main() -> None:
    args = build_argument_parser().parse_args()
    if not TEMPLATE_PPTX.exists():
        raise FileNotFoundError("templates/template.pptx が見つかりません。")

    output_pptx = args.output
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    content = load_source_text(args)
    lines = extract_content_section(content)
    title, slides = parse_slides(lines)

    prs = Presentation(str(TEMPLATE_PPTX))
    # テンプレートに含まれる既存スライドを削除
    while len(prs.slides._sldIdLst):
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    add_title_slide(prs, title)
    for slide_data in slides:
        if slide_data.get("kind") == "section":
            add_section_slide(prs, slide_data)
        else:
            add_content_slide(prs, slide_data)
    prs.save(str(output_pptx))
    if not args.no_open:
        os.startfile(str(output_pptx))


if __name__ == "__main__":
    main()
